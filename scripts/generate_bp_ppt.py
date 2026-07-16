from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUTPUT_PATH = "docs/OraSage_5人最小化运营商业计划书_BP.pptx"


SLIDES = [
    {
        "title": "OraSage 商业计划书（5人最小化运营版）",
        "subtitle": "命理服务 + 电商的一体化商业化方案\n2026-07",
        "bullets": [],
    },
    {
        "title": "项目概览",
        "bullets": [
            "OraSage 是命理服务（bazi/ziwei/tarot）+ 电商（shop）+ 统一账户（auth）平台。",
            "当前 8 个应用架构已基本成型，具备跨应用登录与交易承接能力。",
            "商业化重点：把“可运行”升级为“可持续盈利”。",
        ],
    },
    {
        "title": "市场机会与用户痛点",
        "bullets": [
            "用户需求：情绪压力、决策焦虑、关系与职业困惑持续增长。",
            "现有产品痛点：内容碎片化、可信度不足、服务交付不稳定。",
            "机会窗口：AI + 命理内容 + 电商承接，形成“内容-服务-商品”闭环。",
        ],
    },
    {
        "title": "产品与技术现状",
        "bullets": [
            "已具备：8 应用子域架构、统一登录桥接、商城下单承接。",
            "可支撑：多应用分层部署与跨应用用户识别。",
            "待补齐：真实 VPS 全链路验收、admin/cms 运营化、E2E 与监控。",
        ],
    },
    {
        "title": "商业模式设计",
        "bullets": [
            "收入线 1：数字内容付费（深度报告、AI 解读包）。",
            "收入线 2：会员订阅（月卡/季卡，权益含内容+折扣）。",
            "收入线 3：电商毛利（推荐商品、组合套餐）。",
            "收入线 4：高客单服务（1v1 咨询、年度规划）。",
        ],
    },
    {
        "title": "核心增长漏斗",
        "bullets": [
            "渠道获客：SEO 内容、短视频、社媒矩阵、私域裂变。",
            "转化路径：免费测算 -> 深度报告付费 -> 商品加购 -> 会员复购。",
            "运营重点：A/B 测试（价格、文案、推荐位）、客服 SOP、复购激励。",
        ],
    },
    {
        "title": "5人最小化团队配置",
        "bullets": [
            "负责人（1）：产品与商业策略、关键决策与合作。",
            "技术负责人（1）：稳定性、上线、数据口径与安全。",
            "全栈工程师（1）：支付链路、后台最小功能、实验迭代。",
            "增长运营（1）：内容分发、转化漏斗、用户留存。",
            "内容/交付（1）：命理内容生产、会员服务、客服交付。",
        ],
    },
    {
        "title": "90天执行路线图",
        "bullets": [
            "D1-D30：上线就绪与验收（登录、支付、订单、日志告警）。",
            "D31-D60：商业化首版（分层定价、会员首版、运营流程）。",
            "D61-D90：增长验证（投放小规模复制、复购与留存优化）。",
            "原则：先验证转化，再放大投放。",
        ],
    },
    {
        "title": "产品路线（MVP范围）",
        "bullets": [
            "必做：登录打通、核心测算路径、支付闭环、订单可追踪。",
            "必做：后台最小运营功能（用户/订单/内容）。",
            "延后：复杂社区、重型推荐系统、非关键视觉重构。",
            "目标：以最小功能集形成可验证的营收闭环。",
        ],
    },
    {
        "title": "技术与交付策略",
        "bullets": [
            "架构策略：各应用独立部署，内网 API 调用，边界清晰。",
            "质量策略：关键链路 E2E（登录->测算->下单->支付）。",
            "监控策略：可用性、错误率、支付成功率、核心接口延迟。",
            "安全策略：JWT/密钥管理、支付回调验签、权限最小化。",
        ],
    },
    {
        "title": "运营策略（内容+转化+复购）",
        "bullets": [
            "内容：围绕高需求场景做专题（情感、事业、财富、健康）。",
            "转化：免费内容引导低客单首购，首购后推荐会员。",
            "复购：会员内容日历、节气节点活动、私域触达机制。",
            "服务：标准化咨询模板，缩短交付时间，提升满意度。",
        ],
    },
    {
        "title": "财务模型假设（基准）",
        "bullets": [
            "团队人数：5人；月固定成本（人力+基础运营）：17万~22万。",
            "技术与工具：0.8万~1.5万/月；增长投放：2万~6万/月。",
            "收入结构：数字付费 + 会员订阅 + 电商毛利。",
            "公式：月毛利 = 数字毛利 + 会员毛利 + 电商毛利。",
        ],
    },
    {
        "title": "财务预测（三情景）",
        "bullets": [
            "保守：月毛利 14万，月成本 20万，月经营结果 -6万。",
            "基准：月毛利 24万，月成本 22万，月经营结果 +2万。",
            "进取：月毛利 38万，月成本 26万，月经营结果 +12万。",
            "现金建议：准备 9 个月 runway（180万~220万）。",
        ],
    },
    {
        "title": "关键KPI与里程碑",
        "bullets": [
            "获客：UV、注册转化率、渠道成本。",
            "变现：付费转化率、ARPPU、会员渗透率。",
            "经营：CAC、LTV、回本周期、月毛利率。",
            "稳定性：支付成功率、可用性、P95 延迟。",
        ],
    },
    {
        "title": "风险与融资用途",
        "bullets": [
            "主要风险：支付迁移、系统稳定性、增长效率、合规风险。",
            "应对：灰度上线、监控告警、预算闸门、法务审查。",
            "资金用途：产品迭代、增长投放、内容生产、风险缓冲。",
            "目标：12个月达成可持续增长与月度盈亏平衡。",
        ],
    },
]


def add_slide(prs: Presentation, title: str, bullets: list[str], subtitle: str = "") -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    title_box = slide.shapes.title
    title_box.text = title
    title_tf = title_box.text_frame
    title_tf.paragraphs[0].font.size = Pt(34)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(13, 35, 87)

    body = slide.placeholders[1].text_frame
    body.clear()

    if subtitle:
        p = body.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.bold = False
        p.font.color.rgb = RGBColor(40, 40, 40)
        body.add_paragraph()

    for i, item in enumerate(bullets):
        p = body.paragraphs[0] if (i == 0 and not subtitle) else body.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(32, 32, 32)

    # Footer
    left = Inches(0.5)
    top = Inches(6.9)
    width = Inches(12.3)
    height = Inches(0.3)
    footer = slide.shapes.add_textbox(left, top, width, height)
    footer_tf = footer.text_frame
    footer_tf.text = "OraSage | 5人最小化运营 BP"
    footer_tf.paragraphs[0].font.size = Pt(11)
    footer_tf.paragraphs[0].font.color.rgb = RGBColor(120, 120, 120)


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        add_slide(
            prs,
            title=spec["title"],
            bullets=spec.get("bullets", []),
            subtitle=spec.get("subtitle", ""),
        )

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build_presentation()
    print(f"Generated: {OUTPUT_PATH}")
